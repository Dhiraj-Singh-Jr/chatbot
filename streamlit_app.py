import streamlit as st
import os
import pandas as pd
import fitz  # PyMuPDF
from google import genai
from google.genai import types
from docx import Document
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from bs4 import BeautifulSoup
import re
import base64

# Set Gemini API Key
client = genai.Client(api_key="AIzaSyC_sGvMgnZvJSEg9j1MS0vpBwsgCErLxr0")
os.environ["XDG_CONFIG_HOME"] = "/tmp"
# --------- File Extractors ---------
def extract_text_from_docx(docx_file):
    document = Document(docx_file)
    return "\n".join([para.text for para in document.paragraphs])

def extract_text_from_csv(csv_file):
    df = pd.read_csv(csv_file)
    return df.to_string(index=False)

def extract_text_from_xlsx(xlsx_file):
    df = pd.read_excel(xlsx_file)
    return df.to_string(index=False)

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_text_from_pdf(pdf_file):
    doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_html(html_file):
    soup = BeautifulSoup(html_file.read(), "html.parser")
    return soup.get_text()

def extract_text_from_tex(tex_file):
    content = tex_file.read().decode("utf-8")
    content = re.sub(r'\\[a-zA-Z]+\{[^}]*\}', '', content)
    content = re.sub(r'\\[a-zA-Z]+', '', content)
    return content

def process_image(image_file):
    image_bytes = image_file.read()
    encoded_image = base64.b64encode(image_bytes).decode("utf-8")
    return {
        "inline_data": {
            "mime_type": image_file.type,
            "data": encoded_image
        }
    }

def export_conversation_to_pdf(conversation_history):
    pdf_path = "Conversation.pdf"
    doc = SimpleDocTemplate(pdf_path, pagesize=A4)
    styles = getSampleStyleSheet()
    elements = []

    for i, (user_q, gemini_a) in enumerate(conversation_history):
        elements.append(Paragraph(f"<b>Q{i+1}:</b> {user_q}", styles["Normal"]))
        elements.append(Spacer(1, 8))
        elements.append(Paragraph(f"<b>A{i+1}:</b> {gemini_a.replace(chr(10), '<br/>')}", styles["Normal"]))
        elements.append(Spacer(1, 16))

    doc.build(elements)
    return pdf_path

# --------- Main App ---------
def main():
    st.set_page_config(page_title="Gemini 2.5 Pro Chat Q&A", layout="wide")
    st.title("📄 Chat with Your Documents leveraging Internet Using Gemini 2.5 Pro (Team 18))")

    if "conversation" not in st.session_state:
        st.session_state.conversation = []
    if "documents_text" not in st.session_state:
        st.session_state.documents_text = []
    if "chat_active" not in st.session_state:
        st.session_state.chat_active = True
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []

    uploaded_files = st.file_uploader(
        "Upload files (.txt, .docx, .csv, .xlsx, .pptx, .pdf, .html, .htm, .tex, .jpg, .jpeg, .png):",
        type=["txt", "docx", "csv", "xlsx", "pptx", "pdf", "html", "htm", "tex", "jpg", "jpeg", "png"],
        accept_multiple_files=True
    )

    if uploaded_files:
        all_text_content = ""
        for uploaded_file in uploaded_files:
            filetype = uploaded_file.name.split(".")[-1].lower()
            try:
                if filetype == "txt":
                    all_text_content += uploaded_file.read().decode("utf-8") + "\n"
                elif filetype == "docx":
                    all_text_content += extract_text_from_docx(uploaded_file) + "\n"
                elif filetype == "csv":
                    all_text_content += extract_text_from_csv(uploaded_file) + "\n"
                elif filetype == "xlsx":
                    all_text_content += extract_text_from_xlsx(uploaded_file) + "\n"
                elif filetype == "pptx":
                    all_text_content += extract_text_from_pptx(uploaded_file) + "\n"
                elif filetype == "pdf":
                    all_text_content += extract_text_from_pdf(uploaded_file) + "\n"
                elif filetype in ["html", "htm"]:
                    all_text_content += extract_text_from_html(uploaded_file) + "\n"
                elif filetype == "tex":
                    all_text_content += extract_text_from_tex(uploaded_file) + "\n"
                elif filetype in ["jpg", "jpeg", "png"]:
                    image_data = process_image(uploaded_file)
                    st.session_state.image_data = image_data
                    all_text_content += "Image uploaded and processed.\n"
                    st.image(uploaded_file, caption=uploaded_file.name, use_container_width=True)
                else:
                    st.error(f"Unsupported file format: {uploaded_file.name}")
            except Exception as e:
                st.error(f"Failed to extract text from {uploaded_file.name}: {e}")

        st.session_state.documents_text = all_text_content

    if st.session_state.documents_text:
        st.markdown("### 💬 Conversation")
        for user_q, gemini_a in st.session_state.conversation:
            st.markdown(f"**You:** {user_q}")
            st.markdown(f"**Gemini:**\n\n{gemini_a}")

        if st.session_state.chat_active:
            with st.form(key="chat_form", clear_on_submit=True):
                user_input = st.text_input("Ask a question about the documents (type 'exit' to stop):")
                submit = st.form_submit_button("Send")

                if submit and user_input:
                    if user_input.strip().lower() == "exit":
                        st.session_state.chat_active = False
                        st.success("Chat ended. Reload to start again.")
                    else:
                        with st.spinner("Gemini is thinking..."):
                            content_blocks = []

                            if st.session_state.conversation:
                                history_text = "\n\n".join(
                                    f"Q: {entry[0]}\nA: {entry[1]}"
                                    for entry in st.session_state.conversation
                                )
                                content_blocks.append({"text": f"Previous conversation:\n{history_text}"})

                            if st.session_state.documents_text.strip():
                                content_blocks.append({"text": f"Context:\n{st.session_state.documents_text}"})

                            if "image_data" in st.session_state:
                                content_blocks.append(st.session_state.image_data)

                            content_blocks.append({"text": f"Question: {user_input}"})

                            # ✅ Gemini with Internet Search
                            response = client.models.generate_content(
                                model="gemini-2.5-pro-preview-05-06",
                                contents=content_blocks,
                                config={"tools": [{"google_search": {}}]}
                            )

                            st.session_state.conversation.append((user_input, response.text))
                            st.success("💡 Answer:")
                            st.write(response.text)

                            st.session_state.chat_history.append({
                                "question": user_input,
                                "answer": response.text
                            })

        if st.session_state.conversation:
            pdf_path = export_conversation_to_pdf(st.session_state.conversation)
            with open(pdf_path, "rb") as f:
                st.download_button(
                    label="📥 Export Conversation as PDF",
                    data=f,
                    file_name="Conversation.pdf",
                    mime="application/pdf"
                )

if __name__ == "__main__":
    main()
